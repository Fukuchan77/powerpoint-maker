"""Unit tests for the content extractor service.

Target coverage: 90%
"""

import json
from unittest.mock import MagicMock, patch

import pytest

from app.schemas import AnalysisMode
from app.services.extractor import ContentExtractor


class TestContentExtractor:
    """Tests for the ContentExtractor class."""

    def setup_method(self):
        """Set up test fixtures."""
        self.extractor = ContentExtractor(base_url="http://localhost:8000")

    def test_init_strips_trailing_slash(self):
        """Test that base_url has trailing slash stripped."""
        extractor = ContentExtractor(base_url="http://localhost:8000/")
        assert extractor.base_url == "http://localhost:8000"

    @pytest.mark.asyncio
    async def test_extract_creates_extraction_directory(self, tmp_path, monkeypatch):
        """Test that extraction creates necessary directories."""
        # Mock EXTRACTED_IMAGES_DIR
        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        # Create a minimal mock presentation
        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = []
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.CONTENT)

        # Verify directory was created
        extraction_dir = tmp_path / result.extraction_id
        assert extraction_dir.exists()
        assert (extraction_dir / "images").exists()
        assert (extraction_dir / "metadata.json").exists()

    @pytest.mark.asyncio
    async def test_extract_returns_valid_result(self, tmp_path, monkeypatch):
        """Test that extract returns valid ContentExtractionResult."""
        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = []
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.CONTENT)

        assert result.extraction_id is not None
        assert result.filename == "test.pptx"
        assert result.expires_at is not None
        assert isinstance(result.slides, list)
        assert isinstance(result.images, list)
        assert isinstance(result.warnings, list)

    @pytest.mark.asyncio
    async def test_extract_saves_metadata(self, tmp_path, monkeypatch):
        """Test that metadata is saved correctly."""
        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = []
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.CONTENT)

        metadata_path = tmp_path / result.extraction_id / "metadata.json"
        metadata = json.loads(metadata_path.read_text())

        assert "expires_at" in metadata
        assert "image_count" in metadata
        assert "filename" in metadata
        assert metadata["filename"] == "test.pptx"

    def test_extract_text_from_shape(self):
        """Test text extraction from shape."""
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        mock_para = MagicMock()
        mock_para.text = "  Test text  "
        mock_shape.text_frame.paragraphs = [mock_para]

        texts = self.extractor._extract_text_from_shape(mock_shape)
        assert texts == ["Test text"]

    def test_extract_text_from_shape_empty(self):
        """Test text extraction from shape with no text."""
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        mock_para = MagicMock()
        mock_para.text = "   "
        mock_shape.text_frame.paragraphs = [mock_para]

        texts = self.extractor._extract_text_from_shape(mock_shape)
        assert texts == []

    def test_extract_text_from_shape_no_text_frame(self):
        """Test text extraction from shape without text frame."""
        mock_shape = MagicMock()
        mock_shape.has_text_frame = False

        texts = self.extractor._extract_text_from_shape(mock_shape)
        assert texts == []

    def test_extract_image(self, tmp_path):
        """Test image extraction."""
        extraction_dir = tmp_path / "test_extraction"
        extraction_dir.mkdir()
        (extraction_dir / "images").mkdir()

        mock_shape = MagicMock()
        mock_shape.image.ext = "png"
        mock_shape.image.blob = b"fake image data"
        mock_shape.image.content_type = "image/png"
        mock_shape.name = "test_image.png"

        result = self.extractor._extract_image(
            mock_shape,
            extraction_id="test-id",
            slide_idx=0,
            extraction_dir=extraction_dir,
        )

        assert result is not None
        assert result.slide_index == 0
        assert result.content_type == "image/png"
        assert "test-id" in result.url

    def test_extract_image_failure(self, tmp_path):
        """Test image extraction failure handling."""
        extraction_dir = tmp_path / "test_extraction"
        extraction_dir.mkdir()
        (extraction_dir / "images").mkdir()

        mock_shape = MagicMock()
        mock_shape.image = property(lambda self: (_ for _ in ()).throw(AttributeError("no image")))

        # Simulate AttributeError
        del mock_shape.image

        result = self.extractor._extract_image(
            mock_shape,
            extraction_id="test-id",
            slide_idx=0,
            extraction_dir=extraction_dir,
        )

        assert result is None

    def test_extract_chart(self):
        """Test chart extraction."""
        mock_shape = MagicMock()
        mock_shape.chart.chart_type = "BAR_CLUSTERED"
        mock_shape.chart.plots = [MagicMock()]
        mock_shape.chart.plots[0].categories = ["A", "B", "C"]

        mock_series = MagicMock()
        mock_series.name = "Series 1"
        mock_series.values = [1.0, 2.0, 3.0]
        mock_shape.chart.plots[0].series = [mock_series]

        result = self.extractor._extract_chart(mock_shape, slide_idx=0)

        assert result is not None
        assert result.slide_index == 0
        assert result.chart_type == "BAR_CLUSTERED"
        assert result.categories == ["A", "B", "C"]
        assert len(result.series) == 1

    def test_extract_chart_failure(self):
        """Test chart extraction failure handling."""
        mock_shape = MagicMock()
        mock_shape.chart.plots = []  # Empty plots will cause IndexError

        result = self.extractor._extract_chart(mock_shape, slide_idx=0)

        assert result is None

    def test_handle_unsupported_shape(self):
        """Test warning generation for unsupported shapes."""
        mock_shape = MagicMock()
        mock_shape.shape_type = "VIDEO"

        warning = self.extractor._handle_unsupported_shape(mock_shape, 0, "Video not supported")

        assert "Slide 1" in warning
        assert "VIDEO" in warning
        assert "Video not supported" in warning

    @pytest.mark.asyncio
    async def test_template_mode_skips_content(self, tmp_path, monkeypatch):
        """Test that template mode skips content extraction."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        # Create mock slide with shapes
        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [MagicMock(text="Content")]
        mock_shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        mock_shape.has_chart = False
        mock_slide.shapes = [mock_shape]
        mock_slide.slide_layout.slide_master.slide_layouts.index.return_value = 1

        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = [mock_slide]
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.TEMPLATE)

        # In template mode, body_text should be empty
        assert len(result.slides) == 1
        assert result.slides[0].body_text == []


class TestExtractSlideEdgeCases:
    """Edge case tests for _extract_slide method."""

    def setup_method(self):
        """Set up test fixtures."""
        self.extractor = ContentExtractor(base_url="http://localhost:8000")

    def test_layout_index_fallback_on_value_error(self, tmp_path):
        """Test layout_idx defaults to 0 when ValueError is raised."""
        mock_slide = MagicMock()
        mock_slide.shapes = []
        mock_slide.slide_layout.slide_master.slide_layouts.index.side_effect = ValueError("not in list")

        result, _, _ = self.extractor._extract_slide(mock_slide, 0, "test-id", tmp_path, AnalysisMode.CONTENT)

        assert result.layout_index == 0

    def test_layout_index_fallback_on_attribute_error(self, tmp_path):
        """Test layout_idx defaults to 0 when AttributeError is raised."""
        mock_slide = MagicMock()
        mock_slide.shapes = []
        mock_slide.slide_layout.slide_master = None  # Causes AttributeError on .slide_layouts

        result, _, _ = self.extractor._extract_slide(mock_slide, 0, "test-id", tmp_path, AnalysisMode.CONTENT)

        assert result.layout_index == 0

    def test_shape_exception_adds_warning(self, tmp_path):
        """Test that shape processing exception adds a warning."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = property(lambda self: (_ for _ in ()).throw(RuntimeError("shape error")))
        # Force exception via has_text_frame True but broken paragraphs access
        del mock_shape.text_frame.paragraphs
        mock_shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        mock_shape.has_chart = False
        mock_slide.shapes = [mock_shape]
        mock_slide.slide_layout.slide_master.slide_layouts.index.return_value = 1

        _, _, warnings = self.extractor._extract_slide(mock_slide, 0, "test-id", tmp_path, AnalysisMode.CONTENT)

        # The shape should have caused an exception and added a warning
        assert len(warnings) == 1
        assert "Slide 1" in warnings[0]

    def test_title_extraction_with_placeholder(self, tmp_path):
        """Test title extraction from title placeholder."""
        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True
        mock_shape.placeholder_format.type = 1  # TITLE
        mock_shape.has_text_frame = True
        mock_shape.text_frame.text = "  Test Title  "
        mock_slide.shapes = [mock_shape]
        mock_slide.slide_layout.slide_master.slide_layouts.index.return_value = 0

        result, _, _ = self.extractor._extract_slide(mock_slide, 0, "test-id", tmp_path, AnalysisMode.CONTENT)

        assert result.title == "Test Title"

    def test_image_extraction_with_none_extension(self, tmp_path):
        """Test image extraction when extension is None."""
        extraction_dir = tmp_path / "test_extraction"
        extraction_dir.mkdir()
        (extraction_dir / "images").mkdir()

        mock_shape = MagicMock()
        mock_shape.image.ext = None
        mock_shape.image.blob = b"fake image data"
        mock_shape.image.content_type = None
        mock_shape.name = None

        result = self.extractor._extract_image(
            mock_shape,
            extraction_id="test-id",
            slide_idx=0,
            extraction_dir=extraction_dir,
        )

        assert result is not None
        # Should default to png
        assert result.content_type == "image/png"

    def test_chart_extraction_with_empty_name(self):
        """Test chart extraction when series name is None."""
        mock_shape = MagicMock()
        mock_shape.chart.chart_type = "LINE"
        mock_shape.chart.plots = [MagicMock()]
        mock_shape.chart.plots[0].categories = ["Q1", "Q2"]

        mock_series = MagicMock()
        mock_series.name = None  # Empty name
        mock_series.values = [10.0, 20.0]
        mock_shape.chart.plots[0].series = [mock_series]

        result = self.extractor._extract_chart(mock_shape, slide_idx=0)

        assert result is not None
        assert result.series[0].name == ""

    def test_chart_extraction_with_no_categories(self):
        """Test chart extraction when categories is None."""
        mock_shape = MagicMock()
        mock_shape.chart.chart_type = "PIE"
        mock_shape.chart.plots = [MagicMock()]
        mock_shape.chart.plots[0].categories = None

        mock_series = MagicMock()
        mock_series.name = "Data"
        mock_series.values = [30.0, 70.0]
        mock_shape.chart.plots[0].series = [mock_series]

        result = self.extractor._extract_chart(mock_shape, slide_idx=0)

        assert result is not None
        assert result.categories == []

    @pytest.mark.asyncio
    async def test_extract_multiple_slides_with_content(self, tmp_path, monkeypatch):
        """Test extraction with multiple slides containing various content."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        # Create mock slides
        mock_slide1 = MagicMock()
        mock_text_shape = MagicMock()
        mock_text_shape.is_placeholder = False
        mock_text_shape.has_text_frame = True
        mock_para = MagicMock()
        mock_para.text = "Content text"
        mock_text_shape.text_frame.paragraphs = [mock_para]
        mock_text_shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        mock_text_shape.has_chart = False
        mock_slide1.shapes = [mock_text_shape]
        mock_slide1.slide_layout.slide_master.slide_layouts.index.return_value = 0

        mock_slide2 = MagicMock()
        mock_slide2.shapes = []
        mock_slide2.slide_layout.slide_master.slide_layouts.index.return_value = 1

        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = [mock_slide1, mock_slide2]
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.CONTENT)

        assert len(result.slides) == 2
        assert result.slides[0].body_text == ["Content text"]
        assert result.slides[1].body_text == []

    @pytest.mark.asyncio
    async def test_extract_group_shape_graceful_skip(self, tmp_path, monkeypatch):
        """Test that group shapes are skipped gracefully."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        monkeypatch.setattr("app.services.extractor.EXTRACTED_IMAGES_DIR", tmp_path)

        mock_slide = MagicMock()
        mock_group = MagicMock()
        mock_group.shape_type = MSO_SHAPE_TYPE.GROUP
        mock_group.has_text_frame = False
        mock_group.has_chart = False

        # Ensure it doesn't crash on property access
        mock_slide.shapes = [mock_group]
        mock_slide.slide_layout.slide_master.slide_layouts.index.return_value = 0

        with patch("app.services.extractor.Presentation") as mock_prs:
            mock_prs.return_value.slides = [mock_slide]
            result = await self.extractor.extract(tmp_path / "test.pptx", AnalysisMode.CONTENT)

        # Should extraction successful but empty content
        assert len(result.slides) == 1
        assert result.slides[0].body_text == []

    def test_extract_slide_catches_image_io_error(self, tmp_path):
        """Test that IOError during image extraction is caught at slide level."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        extraction_dir = tmp_path / "test_extraction"

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        mock_shape.image.ext = "png"

        # Setup blob to raise IOError
        type(mock_shape.image).blob = property(fget=MagicMock(side_effect=IOError("Read failed")))

        mock_slide.shapes = [mock_shape]
        mock_slide.slide_layout.slide_master.slide_layouts.index.return_value = 0

        _, _, warnings = self.extractor._extract_slide(mock_slide, 0, "test-id", extraction_dir, AnalysisMode.CONTENT)

        # Should calculate warning
        assert len(warnings) == 1
        assert "Read failed" in warnings[0]
