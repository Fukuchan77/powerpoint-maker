import os

import pytest
from pptx import Presentation

from app.services.template import TemplateAnalyzer


@pytest.fixture
def sample_pptx():
    # Create a dummy PPTX file
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Hello, World!"

    filename = "test_template.pptx"
    prs.save(filename)
    yield filename
    if os.path.exists(filename):
        os.remove(filename)


def test_analyze_template(sample_pptx):
    analyzer = TemplateAnalyzer()
    result = analyzer.analyze(sample_pptx)

    assert result.filename == sample_pptx
    assert len(result.masters) > 0
    # Basic check ensuring we got layouts back
    assert len(result.masters[0].layouts) > 0
