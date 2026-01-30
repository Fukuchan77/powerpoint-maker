from unittest.mock import Mock, patch

import pytest
from pptx.enum.shapes import PP_PLACEHOLDER

from app.services.template import LayoutRegistry, TemplateAnalyzer


@pytest.fixture
def mock_presentation_class():
    with patch("app.services.template.Presentation") as MockPres:
        yield MockPres


def test_analyze_template_structure(mock_presentation_class):
    # Setup mock presentation structure
    mock_prs = mock_presentation_class.return_value
    mock_master = Mock()
    mock_master.name = "Master 1"

    mock_layout = Mock()
    mock_layout.name = "Layout 1"

    # Mock Placeholders
    ph1 = Mock()
    ph1.name = "Title"
    ph1.placeholder_format.type = PP_PLACEHOLDER.TITLE
    ph1.placeholder_format.idx = 0
    ph1.has_text_frame = True
    ph1.width = 100
    ph1.height = 100
    ph1.left = 0
    ph1.top = 0

    mock_layout.placeholders = [ph1]

    mock_master.slide_layouts = [mock_layout]
    mock_prs.slide_masters = [mock_master]

    analyzer = TemplateAnalyzer()
    result = analyzer.analyze("fake_path.pptx", "id_123")

    assert result.template_id == "id_123"
    assert len(result.masters) == 1
    assert result.masters[0].name == "Master 1"
    assert len(result.masters[0].layouts) == 1
    assert result.masters[0].layouts[0].name == "Layout 1"
    assert result.masters[0].layouts[0].placeholders[0].type == str(PP_PLACEHOLDER.TITLE)


def test_layout_registry_singleton_cache(mock_presentation_class):
    # Test Singleton behavior and caching
    reg1 = LayoutRegistry()
    reg2 = LayoutRegistry()

    assert reg1 is reg2

    # Clear cache to start fresh
    reg1.clear()

    # Analyze once
    mock_presentation_class.return_value.slide_masters = []  # simple empty

    result1 = reg1.get_or_analyze("path1.pptx", "id_1")
    # assert "id_1" in reg1._cache  # No longer accessing _cache directly

    # Analyze again (should hit cache)
    result2 = reg2.get_or_analyze("path1.pptx", "id_1")  # same ID

    assert result1 is result2
    # Verify Presentation initialized only once for id_1
    # Note: mock_presentation_class is the class, not the instance.
    assert mock_presentation_class.call_count == 1

    # Different ID -> analyze called
    reg1.get_or_analyze("path2.pptx", "id_2")
    assert mock_presentation_class.call_count == 2
