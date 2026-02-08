"""Unit tests for Two-Column layout support in PresentationGenerator.

Tests the new Two-Column bullet handling logic added to generator.py,
focusing on the _find_all_body_placeholders method and bullet distribution logic.
"""

from unittest.mock import Mock

import pytest
from pptx.enum.shapes import PP_PLACEHOLDER

from app.services.generator import PresentationGenerator


class TestFindAllBodyPlaceholders:
    """Test the _find_all_body_placeholders method."""

    @pytest.fixture
    def generator(self):
        """Create a PresentationGenerator instance."""
        return PresentationGenerator()

    def create_mock_placeholder(self, ph_type, has_text_frame=True):
        """Helper to create a mock placeholder."""
        ph = Mock()
        ph.placeholder_format = Mock()
        ph.placeholder_format.type = ph_type
        ph.has_text_frame = has_text_frame
        return ph

    def test_find_two_body_placeholders(self, generator):
        """Test finding two BODY placeholders on a slide."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        slide.placeholders = [ph1, ph2]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 2
        assert result[0] == ph1
        assert result[1] == ph2

    def test_find_body_and_object_placeholders(self, generator):
        """Test finding BODY and OBJECT placeholders."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.OBJECT)
        ph3 = self.create_mock_placeholder(PP_PLACEHOLDER.TITLE)
        slide.placeholders = [ph1, ph2, ph3]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 2
        assert result[0] == ph1
        assert result[1] == ph2

    def test_filters_placeholders_without_text_frame(self, generator):
        """Test that placeholders without text_frame are filtered out."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY, has_text_frame=True)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY, has_text_frame=False)
        ph3 = self.create_mock_placeholder(PP_PLACEHOLDER.OBJECT, has_text_frame=True)
        slide.placeholders = [ph1, ph2, ph3]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 2
        assert result[0] == ph1
        assert result[1] == ph3

    def test_returns_empty_list_when_no_body_placeholders(self, generator):
        """Test when no BODY placeholders exist."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.TITLE)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.PICTURE)
        slide.placeholders = [ph1, ph2]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 0

    def test_preserves_placeholder_order(self, generator):
        """Test that placeholders are returned in the order they appear."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.OBJECT)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        ph3 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        slide.placeholders = [ph1, ph2, ph3]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 3
        assert result[0] == ph1  # OBJECT comes first
        assert result[1] == ph2  # First BODY
        assert result[2] == ph3  # Second BODY

    def test_handles_multiple_object_placeholders(self, generator):
        """Test finding multiple OBJECT placeholders."""
        slide = Mock()
        ph1 = self.create_mock_placeholder(PP_PLACEHOLDER.OBJECT)
        ph2 = self.create_mock_placeholder(PP_PLACEHOLDER.OBJECT)
        slide.placeholders = [ph1, ph2]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 2
        assert result[0] == ph1
        assert result[1] == ph2

    def test_handles_empty_placeholders_list(self, generator):
        """Test when slide has no placeholders at all."""
        slide = Mock()
        slide.placeholders = []

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 0

    def test_handles_mixed_placeholder_types(self, generator):
        """Test with a realistic mix of placeholder types."""
        slide = Mock()
        ph_title = self.create_mock_placeholder(PP_PLACEHOLDER.TITLE)
        ph_body1 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        ph_picture = self.create_mock_placeholder(PP_PLACEHOLDER.PICTURE)
        ph_body2 = self.create_mock_placeholder(PP_PLACEHOLDER.BODY)
        ph_chart = self.create_mock_placeholder(PP_PLACEHOLDER.CHART)
        slide.placeholders = [ph_title, ph_body1, ph_picture, ph_body2, ph_chart]

        result = generator._find_all_body_placeholders(slide)

        assert len(result) == 2
        assert result[0] == ph_body1
        assert result[1] == ph_body2


class TestTwoColumnLogicIntegration:
    """Integration tests for Two-Column bullet handling logic.

    These tests verify the logic flow without full generator mocking.
    """

    def test_two_column_detection_with_bullets_right(self):
        """Test that bullets_right field triggers Two-Column logic."""
        from app.schemas import BulletPoint, SlideContent

        # Create slide content with bullets_right (Two-Column indicator)
        slide_content = SlideContent(
            title="Test",
            layout_index=0,
            bullets=[BulletPoint(text="Left", level=0)],
            bullets_right=[BulletPoint(text="Right", level=0)],
        )

        # Verify bullets_right is set
        assert slide_content.bullets_right is not None
        assert len(slide_content.bullets_right) == 1
        assert slide_content.bullets_right[0].text == "Right"

    def test_single_column_without_bullets_right(self):
        """Test that absence of bullets_right indicates single-column layout."""
        from app.schemas import BulletPoint, SlideContent

        # Create slide content without bullets_right (single-column)
        slide_content = SlideContent(title="Test", layout_index=0, bullets=[BulletPoint(text="Bullet", level=0)])

        # Verify bullets_right is None (default)
        assert slide_content.bullets_right is None

    def test_backward_compatibility_with_bullet_points(self):
        """Test that old bullet_points field still works (backward compatibility)."""
        from app.schemas import SlideContent

        # Create slide content using old bullet_points field
        slide_content = SlideContent(title="Test", layout_index=0, bullet_points=["Bullet 1", "Bullet 2"])

        # Verify bullet_points is set
        assert len(slide_content.bullet_points) == 2
        assert slide_content.bullets is None
        assert slide_content.bullets_right is None


# Made with Bob
