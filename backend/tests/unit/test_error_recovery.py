from unittest.mock import Mock

from pptx.enum.shapes import PP_PLACEHOLDER

from app.services.generator import SlidePopulator


def test_safe_get_placeholder_fallback():
    """Test safe_get_placeholder attempts fallback index"""
    mock_slide = Mock()
    mock_slide.placeholders = {}  # Empty placeholders dict-like access

    populator = SlidePopulator(mock_slide)

    # Mock behavior: primary index raises KeyError, fallback doesn't
    # But wait, python-pptx placeholders is an object acting like a dict but also iterable?
    # Usually placeholders[idx] works.

    ph_fallback = Mock()
    ph_fallback.placeholder_format.idx = 5

    # Custom Mock that behaves like pptx Placeholders
    class MockPlaceholders:
        def __init__(self, items):
            self._items = items  # list of placeholders

        def __getitem__(self, key):
            # Dict-like lookup by idx
            for item in self._items:
                if item.placeholder_format.idx == key:
                    return item
            raise KeyError(key)

        def __iter__(self):
            return iter(self._items)

    mock_slide.placeholders = MockPlaceholders([ph_fallback])

    # 1. Primary exists (mock setup required differently if we want primary to exist)
    # Let's test fallback scenario
    result = populator.safe_get_placeholder(10, fallback_idx=5)
    assert result == ph_fallback

    # 2. Both missing
    result = populator.safe_get_placeholder(10, fallback_idx=99)
    assert result is None
    assert "fallback=99" in populator.errors[-1]


def test_validate_content_type():
    """Test validate_content_type checks"""
    populator = SlidePopulator(Mock())

    # Text placeholder
    ph_text = Mock()
    ph_text.placeholder_format.type = PP_PLACEHOLDER.BODY

    assert populator.validate_content_type(ph_text, "text") is True
    assert populator.validate_content_type(ph_text, "image") is True  # Dictionary said BODY accepts image?
    # Let's check logic:
    # "image": [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY],
    # Yes.

    ph_pic = Mock()
    ph_pic.placeholder_format.type = PP_PLACEHOLDER.PICTURE

    assert populator.validate_content_type(ph_pic, "image") is True
    assert populator.validate_content_type(ph_pic, "text") is False

    # None placeholder
    assert populator.validate_content_type(None, "text") is False
