import os
from io import BytesIO
from unittest.mock import Mock, patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.schemas import SlideContent
from app.services.generator import PresentationGenerator, SlidePopulator

# Note: sample_pptx fixture is provided by conftest.py


def test_generate_presentation(sample_pptx, tmp_path):
    generator = PresentationGenerator()

    slides = [
        SlideContent(
            layout_index=1,  # Title and Content usually
            title="Slide 1",
            bullet_points=["Point A", "Point B"],
        ),
        SlideContent(layout_index=1, title="Slide 2", bullet_points=["Point C"]),
    ]

    output = str(tmp_path / "output_gen.pptx")
    generated_path = generator.generate(sample_pptx, slides, output)

    assert os.path.exists(generated_path)

    # Verify content
    prsk = Presentation(generated_path)
    assert len(prsk.slides) == 2
    assert prsk.slides[0].shapes.title.text == "Slide 1"


def test_generate_invalid_template_path():
    generator = PresentationGenerator()
    with pytest.raises(FileNotFoundError):
        generator.generate("non_existent.pptx", [], "out.pptx")


@patch("builtins.print")
def test_generate_invalid_layout_index(mock_print, sample_pptx, tmp_path):
    generator = PresentationGenerator()
    slides = [SlideContent(layout_index=99, title="T", bullet_points=[])]

    out = str(tmp_path / "out_warn.pptx")
    generator.generate(sample_pptx, slides, out)

    # Check that it printed a warning
    args, _ = mock_print.call_args
    assert "Warning" in args[0]

    # Check correct completion (file exists, but maybe empty content for that slide)
    assert os.path.exists(out)


@pytest.mark.parametrize("status_code", [404, 500])
def test_generate_image_failure(sample_pptx, tmp_path, status_code):
    generator = PresentationGenerator()
    slides = [SlideContent(layout_index=1, title="Img", bullet_points=[], image_url="http://fail.com/img.png")]

    # We need to ensure logic doesn't crash.
    # We can mock requests.get
    with patch("requests.get") as mock_get:
        mock_get.return_value = Mock(status_code=status_code)

        out = str(tmp_path / "out_img_fail.pptx")
        generator.generate(sample_pptx, slides, out)

        assert os.path.exists(out)
        prs = Presentation(out)
        # Check that slide was created (even if image failed)
        assert len(prs.slides) == 1


# --- SlidePopulator Tests ---


def test_safe_get_placeholder_success():
    mock_slide = Mock()
    mock_ph = Mock()
    mock_slide.placeholders = {1: mock_ph}

    populator = SlidePopulator(mock_slide)
    assert populator.safe_get_placeholder(1) == mock_ph


def test_safe_get_placeholder_fallback():
    mock_slide = Mock()
    mock_ph_correct = Mock()
    mock_ph_correct.placeholder_format.idx = 2

    # Simulate dict lookup failure but list presence
    mock_placeholders = Mock()
    mock_placeholders.__getitem__ = Mock(side_effect=KeyError)
    mock_placeholders.__iter__ = Mock(return_value=iter([mock_ph_correct]))

    mock_slide.placeholders = mock_placeholders

    populator = SlidePopulator(mock_slide)
    assert populator.safe_get_placeholder(2) == mock_ph_correct


def test_safe_get_placeholder_failure():
    mock_slide = Mock()
    mock_placeholders = Mock()
    mock_placeholders.__getitem__ = Mock(side_effect=KeyError)
    mock_placeholders.__iter__ = Mock(return_value=iter([]))  # Empty

    mock_slide.placeholders = mock_placeholders

    populator = SlidePopulator(mock_slide)
    result = populator.safe_get_placeholder(99)
    assert result is None
    assert "No placeholder with idx=99 (fallback=None)" in populator.errors


def test_replace_text_preserve_format_no_runs():
    mock_p = Mock()
    mock_p.runs = []

    populator = SlidePopulator(None)
    populator.replace_text_preserve_format(mock_p, "New Text")

    mock_p.add_run.assert_called_once()
    mock_p.add_run().text = "New Text"


def test_replace_text_preserve_format_existing_runs():
    # Setup a paragraph with multiple runs
    # Structure: Paragraph has .runs list, and ._p element which has children

    # We need to mock the underlying lxml structure slightly because the code iterates over p._p

    mock_p_element = Mock()

    # Create mock elements for runs
    r1 = Mock()
    r1.tag = "...a:r"  # Ends with r
    r2 = Mock()
    r2.tag = "...a:r"

    # And maybe a non-run element
    other = Mock()
    other.tag = "...a:pPr"

    # list(p) will return these
    mock_p_element.__iter__ = Mock(return_value=iter([other, r1, r2]))

    mock_paragraph = Mock()
    mock_paragraph._p = mock_p_element

    # .runs list usually filters for 'r' elements.
    # The code accesses .runs[0] to set text.
    mock_run_1 = Mock()
    mock_paragraph.runs = [mock_run_1, Mock()]  # 2 runs initially

    populator = SlidePopulator(None)
    populator.replace_text_preserve_format(mock_paragraph, "Updated")

    # Verification
    # 1. remove should be called for r2 (the second run)
    mock_p_element.remove.assert_called_once_with(r2)

    # 2. Text updated on first run
    assert mock_run_1.text == "Updated"


def test_insert_picture_fit_success():
    mock_placeholder = Mock()
    mock_picture = Mock()

    # Assume 100x100 placeholder
    mock_picture.width = 1000
    mock_picture.height = 1000
    mock_placeholder.insert_picture.return_value = mock_picture

    # Create valid image bytes (1x1 pixel)
    img = Image.new("RGB", (50, 50), color="red")
    img_byte_arr = BytesIO()
    img.save(img_byte_arr, format="PNG")
    img_bytes = img_byte_arr.getvalue()

    populator = SlidePopulator(None)
    result = populator.insert_picture_fit(mock_placeholder, img_bytes)

    assert result == mock_picture
    assert mock_picture.crop_top == 0


def test_insert_picture_fit_resize_width():
    # Placeholder is wide (200x100, ratio 2.0)
    # Image is square (50x50, ratio 1.0)
    # ph_ratio (2.0) > img_ratio (1.0) -> width = img_ratio * height

    mock_placeholder = Mock()
    mock_picture = Mock()
    mock_picture.width = 200
    mock_picture.height = 100

    mock_placeholder.insert_picture.return_value = mock_picture

    img = Image.new("RGB", (50, 50), color="red")
    img_byte_arr = BytesIO()
    img.save(img_byte_arr, format="PNG")
    img_bytes = img_byte_arr.getvalue()

    populator = SlidePopulator(None)
    populator.insert_picture_fit(mock_placeholder, img_bytes)

    # Expectation: width becomes 1.0 * 100 = 100
    assert mock_picture.width == 100


def test_insert_picture_fit_failure():
    mock_placeholder = Mock()
    mock_placeholder.insert_picture.side_effect = Exception("Insert failed")

    populator = SlidePopulator(None)
    result = populator.insert_picture_fit(mock_placeholder, b"fake_data")

    assert result is None
    assert len(populator.errors) > 0
    assert "Image processing failed" in populator.errors[0]
    assert result is None
    assert len(populator.errors) > 0
    assert "Image processing failed" in populator.errors[0]


def test_populate_bullets_mixed_content():
    # Test mixture of string and BulletPoint objects
    mock_ph = Mock()
    mock_tf = Mock()
    mock_ph.text_frame = mock_tf

    # Setup add_paragraph to return a mock paragraph each time
    p1 = Mock()
    p1.runs = []
    p2 = Mock()
    p2.runs = [Mock()]  # Has a run for styling

    mock_tf.add_paragraph.side_effect = [p1, p2]

    from app.schemas import BulletPoint

    bullets = ["Simple String", BulletPoint(text="Level 1 Item", level=1)]

    populator = SlidePopulator(None)
    populator.populate_bullets(mock_ph, bullets, theme_color="ACCENT_2")

    # Verification
    assert mock_tf.clear.called
    assert mock_tf.add_paragraph.call_count == 2

    # First item (String)
    assert p1.text == "Simple String"
    assert p1.level == 0

    # Second item (BulletPoint)
    assert p2.text == "Level 1 Item"
    assert p2.level == 1

    # Theme Color Check
    # p2 has runs, so it should have attempted to set color
    # We can check if populator.set_theme_color logic ran
    # Ideally checking side effects on run font color
    assert p2.runs[0].font.color.theme_color is not None  # Logic sets it


def test_insert_chart():
    slide_mock = Mock()
    populator = SlidePopulator(slide_mock)

    # Mock placeholder
    placeholder = Mock()
    placeholder.insert_chart = Mock()
    chart_obj = Mock()
    placeholder.insert_chart.return_value = chart_obj

    # Chart Data
    from pptx.enum.chart import XL_CHART_TYPE

    from app.schemas import ChartData, ChartSeries

    chart_data = ChartData(
        title="Sales",
        categories=["A", "B"],
        series=[ChartSeries(name="Series 1", values=[1.0, 2.0])],
        type="COLUMN_CLUSTERED",
    )

    populator.insert_chart(placeholder, chart_data)

    # Check call
    placeholder.insert_chart.assert_called_once()
    args = placeholder.insert_chart.call_args
    assert args[0][0] == XL_CHART_TYPE.COLUMN_CLUSTERED

    # Check title set
    assert chart_obj.chart.has_title is True
    assert chart_obj.chart.chart_title.text_frame.text == "Sales"


@patch("app.services.generator.Presentation")
def test_generate_with_chart(mock_presentation):
    # Setup mocks
    prs = mock_presentation.return_value
    slide = Mock()
    prs.slides.add_slide.return_value = slide
    slide.shapes.title = None  # Avoid title logic failure

    # Setup master/layout
    master = Mock()
    layout = Mock()
    master.slide_layouts = [layout]
    prs.slide_masters = [master]

    # Setup placeholder for chart
    ph_chart = Mock()
    ph_chart.placeholder_format.type = PP_PLACEHOLDER.CHART
    slide.placeholders = [ph_chart]

    # Input Data
    from app.schemas import ChartData, ChartSeries, SlideContent

    slides_content = [
        SlideContent(
            layout_index=0,
            title="Chart Slide",
            chart=ChartData(
                title="Test Chart",
                categories=["X", "Y"],
                series=[ChartSeries(name="S1", values=[10.0, 20.0])],
                type="PIE",
            ),
        )
    ]

    gen = PresentationGenerator()
    with patch("os.path.exists", return_value=True):
        gen.generate("dummy_template.pptx", slides_content, "output.pptx")

    # Verify chart inserted
    ph_chart.insert_chart.assert_called()


@patch("app.services.generator.Presentation")
def test_generate_fallback_scenarios(mock_presentation):
    # Setup - a layout with NO Picture placeholder, but a BODY placeholder
    prs = mock_presentation.return_value
    slide = Mock()
    slide.shapes.title = None

    # Mock Master and Layouts
    master = Mock()
    layout = Mock()
    master.slide_layouts = [layout]
    prs.slide_masters = [master]

    # Return slide when added
    prs.slides.add_slide.return_value = slide

    # Placeholder: type=BODY
    ph_body = Mock()
    ph_body.placeholder_format.type = PP_PLACEHOLDER.BODY
    ph_body.has_text_frame = True  # For bullets fallback check if needed

    slide.placeholders = [ph_body]

    # 1. Test Image Fallback to BODY
    from app.schemas import SlideContent

    slides_content = [
        SlideContent(layout_index=0, title="Fallback", bullet_points=[], image_url="http://valid.com/img.png")
    ]

    gen = PresentationGenerator()

    # Mock requests and insert_picture
    with (
        patch("os.path.exists", return_value=True),
        patch("requests.get") as mock_get,
        patch("builtins.print") as mock_print,
    ):
        mock_get.return_value = Mock(status_code=200, content=b"fakeimg")
        # We need populator to support insert_picture on this mock
        ph_body.insert_picture = Mock()
        ph_body.insert_picture.return_value = Mock(width=100, height=100, crop_top=0)  # Mock picture object

        # We also need Image.open to work
        with patch("PIL.Image.open") as mock_img_open:
            mock_img_open.return_value = Mock(size=(100, 100))

            gen.generate("dummy.pptx", slides_content, "out.pptx")

            # Verify body placeholder was used for image
            if not ph_body.insert_picture.called:
                pytest.fail(f"insert_picture not called. Debug prints: {mock_print.call_args_list}")

            ph_body.insert_picture.assert_called()


def test_japanese_font_support():
    slide = Mock()
    populator = SlidePopulator(slide)

    # Test _contains_japanese
    assert populator._contains_japanese("Hello") is False
    assert populator._contains_japanese("こんにちは") is True
    assert populator._contains_japanese("Text with 漢字") is True

    # Test replace_text_preserve_format
    p = Mock()
    run = Mock()
    p.runs = [run]
    p._p = []  # Empty list for iteration

    populator.replace_text_preserve_format(p, "こんにちは")
    assert run.text == "こんにちは"
    assert run.font.name == "Meiryo UI"

    # Test populate_bullets
    ph = Mock()
    p_bull = Mock()
    run_bull = Mock()
    p_bull.runs = [run_bull]
    ph.text_frame.add_paragraph.return_value = p_bull

    populator.populate_bullets(ph, ["日本語の弾丸"])

    assert run_bull.font.name == "Meiryo UI"


def test_set_japanese_font_error_handling():
    """Test error handling when setting Japanese font fails"""
    slide = Mock()
    populator = SlidePopulator(slide)

    # Create a run that raises an exception when accessing _r
    run = Mock()
    run._r.get_or_add_rPr.side_effect = Exception("XML manipulation failed")

    # Should not crash, but add error to errors list
    populator.set_japanese_font(run)

    assert len(populator.errors) > 0
    assert "Failed to set Japanese font" in populator.errors[0]


def test_replace_text_preserve_format_japanese_no_runs():
    """Test Japanese text replacement when paragraph has no runs"""
    mock_p = Mock()
    mock_p.runs = []
    mock_run = Mock()
    mock_p.add_run.return_value = mock_run

    populator = SlidePopulator(None)
    populator.replace_text_preserve_format(mock_p, "日本語")

    mock_p.add_run.assert_called_once()
    assert mock_run.text == "日本語"
    # Should have attempted to set Japanese font
    assert mock_run.font.name == "Meiryo UI"


def test_set_theme_color_invalid_color():
    """Test theme color setting with invalid color name"""
    slide = Mock()
    populator = SlidePopulator(slide)

    run = Mock()
    # Test with non-existent color name - should use ACCENT_1 as fallback
    populator.set_theme_color(run, "INVALID_COLOR_NAME")

    # Should not crash and should have set some color
    assert run.font.color.theme_color is not None


def test_set_theme_color_error_handling():
    """Test error handling when setting theme color fails"""
    slide = Mock()
    populator = SlidePopulator(slide)

    run = Mock()
    # Make the entire font.color access raise an exception
    type(run.font).color = property(lambda self: (_ for _ in ()).throw(Exception("Color setting failed")))

    # Should catch exception and add to errors
    populator.set_theme_color(run, "ACCENT_1")

    assert len(populator.errors) > 0
    assert "Failed to set theme color" in populator.errors[0]


def test_populate_bullets_unknown_item_type():
    """Test populate_bullets with unknown item type (not string or BulletPoint)"""
    mock_ph = Mock()
    mock_tf = Mock()
    mock_ph.text_frame = mock_tf

    # Create an unknown type (e.g., a number)
    bullets = [123, {"invalid": "dict"}]

    populator = SlidePopulator(None)
    populator.populate_bullets(mock_ph, bullets, theme_color=None)

    # Should clear and not crash
    assert mock_tf.clear.called
    # Paragraphs are added but text is not set for unknown types (they are skipped)
    # The code still calls add_paragraph but doesn't set text on unknown types
    assert mock_tf.add_paragraph.call_count == 2


def test_insert_chart_title_error():
    """Test chart insertion when setting title fails"""
    slide_mock = Mock()
    populator = SlidePopulator(slide_mock)

    placeholder = Mock()
    chart_obj = Mock()
    # Make chart title setting raise an exception
    chart_obj.chart.chart_title.text_frame.text = Mock(side_effect=Exception("Title not supported"))
    placeholder.insert_chart.return_value = chart_obj

    from app.schemas import ChartData, ChartSeries

    chart_data = ChartData(
        title="Test Chart",
        categories=["A", "B"],
        series=[ChartSeries(name="Series 1", values=[1.0, 2.0])],
        type="PIE",
    )

    # Should not crash even if title setting fails
    populator.insert_chart(placeholder, chart_data)

    # Chart should still be inserted
    placeholder.insert_chart.assert_called_once()


def test_insert_chart_complete_failure():
    """Test chart insertion when entire operation fails"""
    slide_mock = Mock()
    populator = SlidePopulator(slide_mock)

    placeholder = Mock()
    placeholder.insert_chart.side_effect = Exception("Chart insertion failed")

    from app.schemas import ChartData, ChartSeries

    chart_data = ChartData(
        title="Test Chart",
        categories=["A"],
        series=[ChartSeries(name="S1", values=[1.0])],
        type="COLUMN_CLUSTERED",
    )

    populator.insert_chart(placeholder, chart_data)

    assert len(populator.errors) > 0
    assert "Failed to insert chart" in populator.errors[0]


@patch("builtins.print")
def test_generate_invalid_image_url(mock_print, sample_pptx, tmp_path):
    """Test generation with invalid image URL (not http/https)"""
    generator = PresentationGenerator()
    slides = [
        SlideContent(
            layout_index=1,
            title="Invalid Image",
            bullet_points=[],
            image_url="ftp://invalid.com/image.png",  # Invalid protocol
        )
    ]

    out = str(tmp_path / "out_invalid_url.pptx")
    generator.generate(sample_pptx, slides, out)

    # Should complete without crashing
    assert os.path.exists(out)

    # Should have printed warning about invalid URL
    print_calls = [str(call) for call in mock_print.call_args_list]
    assert any("Invalid image URL" in call for call in print_calls)


@patch("builtins.print")
@patch("requests.get")
def test_generate_image_fetch_exception(mock_get, mock_print, sample_pptx, tmp_path):
    """Test generation when image fetch raises exception"""
    mock_get.side_effect = Exception("Connection timeout")

    generator = PresentationGenerator()
    slides = [
        SlideContent(
            layout_index=1,
            title="Image Fetch Error",
            bullet_points=[],
            image_url="http://example.com/image.png",
        )
    ]

    out = str(tmp_path / "out_fetch_error.pptx")
    generator.generate(sample_pptx, slides, out)

    # Should complete without crashing
    assert os.path.exists(out)

    # Should have printed error message
    print_calls = [str(call) for call in mock_print.call_args_list]
    assert any("Failed to fetch/insert image" in call for call in print_calls)


@patch("builtins.print")
def test_generate_chart_no_placeholder(mock_print, sample_pptx, tmp_path):
    """Test generation when no suitable chart placeholder is found"""
    from app.schemas import ChartData, ChartSeries

    generator = PresentationGenerator()
    slides = [
        SlideContent(
            layout_index=0,  # Layout with no chart placeholder
            title="Chart Slide",
            bullet_points=[],
            chart=ChartData(
                title="Test Chart",
                categories=["X", "Y"],
                series=[ChartSeries(name="S1", values=[10.0, 20.0])],
                type="COLUMN_CLUSTERED",
            ),
        )
    ]

    out = str(tmp_path / "out_no_chart_ph.pptx")
    generator.generate(sample_pptx, slides, out)

    # Should complete without crashing
    assert os.path.exists(out)

    # Should have printed warning
    print_calls = [str(call) for call in mock_print.call_args_list]
    assert any("No suitable placeholder found for chart" in call for call in print_calls)


def test_validate_content_type():
    """Test content type validation for placeholders"""
    slide = Mock()
    populator = SlidePopulator(slide)

    # Test with None placeholder
    assert populator.validate_content_type(None, "text") is False

    # Test with valid text placeholder
    ph = Mock()
    ph.placeholder_format.type = PP_PLACEHOLDER.BODY
    assert populator.validate_content_type(ph, "text") is True

    # Test with invalid combination
    ph.placeholder_format.type = PP_PLACEHOLDER.PICTURE
    assert populator.validate_content_type(ph, "text") is False

    # Test with image placeholder
    assert populator.validate_content_type(ph, "image") is True
