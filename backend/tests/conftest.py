"""
Shared pytest fixtures for backend tests.

This module consolidates common fixtures used across unit and integration tests
to reduce duplication and improve test maintainability.
"""

from io import BytesIO
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# =============================================================================
# Rate Limiter Reset Fixture
# =============================================================================


@pytest.fixture(autouse=True)
def reset_rate_limiter():
    """
    Reset rate limiter between tests to prevent rate limit errors.

    This fixture runs automatically before each test to clear the rate limiter's
    internal state, ensuring tests don't interfere with each other.
    """
    from app.middleware.rate_limit import limiter

    # Reset the limiter's storage before each test
    limiter.reset()
    yield
    # Clean up after test
    limiter.reset()


# =============================================================================
# File-Based Fixtures
# =============================================================================


@pytest.fixture
def sample_pptx(tmp_path):
    """
    Create a minimal PPTX file for testing.

    Returns the path to a temporary PPTX file that is automatically
    cleaned up after the test.
    """
    filename = tmp_path / "test_template.pptx"
    prs = Presentation()
    prs.save(str(filename))
    return str(filename)


@pytest.fixture
def sample_pptx_with_content(tmp_path):
    """
    Create a PPTX file with a title slide for testing.

    Returns the path to a temporary PPTX file with actual content.
    """
    filename = tmp_path / "test_template_with_content.pptx"
    prs = Presentation()
    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title:
        title.text = "Test Presentation"
    prs.save(str(filename))
    return str(filename)


@pytest.fixture
def output_pptx_path(tmp_path):
    """
    Provide a path for generated presentation output.

    Returns a path in the temporary directory for output files.
    """
    return str(tmp_path / "output.pptx")


# =============================================================================
# Mock Upload File Fixtures
# =============================================================================


def create_mock_upload_file(filename: str, content: bytes, content_type: str):
    """
    Factory function to create mock UploadFile objects.

    Args:
        filename: The name of the uploaded file
        content: The file content as bytes
        content_type: The MIME type of the file

    Returns:
        A MagicMock object that behaves like FastAPI's UploadFile
    """
    file = MagicMock()
    file.filename = filename
    file.content_type = content_type
    file.file = BytesIO(content)
    file.size = len(content)

    async def async_read():
        return content

    async def async_seek(pos):
        return file.file.seek(pos)

    file.read = async_read
    file.seek = async_seek
    return file


@pytest.fixture
def mock_upload_file_factory():
    """
    Provide the mock upload file factory function as a fixture.

    Usage:
        def test_something(mock_upload_file_factory):
            file = mock_upload_file_factory("test.pptx", b"content", "application/...")
    """
    return create_mock_upload_file


@pytest.fixture
def valid_pptx_upload():
    """
    Create a mock UploadFile with valid PPTX magic bytes.

    Returns a mock file that passes PPTX validation.
    """
    # PPTX files are ZIP archives, magic bytes: PK\x03\x04
    content = b"PK\x03\x04" + b"\x00" * 1000
    return create_mock_upload_file(
        "valid_template.pptx",
        content,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@pytest.fixture
def invalid_pptx_upload():
    """
    Create a mock UploadFile with invalid magic bytes.

    Returns a mock file that fails PPTX validation.
    """
    content = b"INVALID_CONTENT" + b"\x00" * 100
    return create_mock_upload_file(
        "invalid.pptx",
        content,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# =============================================================================
# Mock Presentation Fixtures
# =============================================================================


@pytest.fixture
def mock_presentation_class():
    """
    Mock the python-pptx Presentation class.

    Useful for testing template analysis without actual PPTX files.
    """
    with patch("app.services.template.Presentation") as MockPres:
        yield MockPres


@pytest.fixture
def mock_placeholder():
    """
    Create a mock placeholder shape.

    Returns a MagicMock configured like a PPTX placeholder.
    """
    placeholder = MagicMock()
    placeholder.name = "Title"
    placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    placeholder.placeholder_format.idx = 0
    placeholder.has_text_frame = True
    placeholder.width = 100
    placeholder.height = 100
    placeholder.left = 0
    placeholder.top = 0
    return placeholder


@pytest.fixture
def mock_slide_layout(mock_placeholder):
    """
    Create a mock slide layout with placeholders.

    Returns a MagicMock configured like a PPTX slide layout.
    """
    layout = MagicMock()
    layout.name = "Title Slide"
    layout.placeholders = [mock_placeholder]
    return layout


@pytest.fixture
def mock_slide_master(mock_slide_layout):
    """
    Create a mock slide master with layouts.

    Returns a MagicMock configured like a PPTX slide master.
    """
    master = MagicMock()
    master.name = "Office Theme"
    master.slide_layouts = [mock_slide_layout]
    return master


# =============================================================================
# Mock LLM Fixtures
# =============================================================================


@pytest.fixture
def mock_llm_response():
    """
    Factory for creating mock LLM responses with JSON content.

    Usage:
        def test_something(mock_llm_response):
            response = mock_llm_response('{"key": "value"}')
    """

    def _create_response(json_content: str):
        mock_response = MagicMock()
        mock_response.state.message.content = f"```json\n{json_content}\n```"
        return mock_response

    return _create_response


@pytest.fixture
def mock_llm():
    """
    Create a mock LLM instance.

    Returns a MagicMock with an async run method.
    """
    mock = MagicMock()
    mock.run = AsyncMock()
    return mock


# =============================================================================
# Research Agent Fixtures
# =============================================================================


@pytest.fixture
def mock_research_dependencies():
    """
    Mock all research agent dependencies.

    Patches DuckDuckGoSearchTool, get_llm, and DocumentConverter.
    Yields a dict with all mock objects for test configuration.
    """
    with (
        patch("app.services.research.DuckDuckGoSearchTool") as MockSearch,
        patch("app.services.research.get_llm") as mock_get_llm,
        patch("app.services.research.DocumentConverter") as MockConverter,
    ):
        # Configure search tool mock
        mock_search_instance = MockSearch.return_value
        mock_search_instance.run = AsyncMock(return_value="Search results")

        # Configure LLM mock
        mock_llm_instance = MagicMock()
        mock_llm_instance.run = AsyncMock()
        mock_get_llm.return_value = mock_llm_instance

        # Configure document converter mock
        mock_converter_instance = MockConverter.return_value
        mock_converter_instance.convert = MagicMock(return_value=MagicMock(document=""))

        yield {
            "search_class": MockSearch,
            "search_instance": mock_search_instance,
            "get_llm": mock_get_llm,
            "llm_instance": mock_llm_instance,
            "converter_class": MockConverter,
            "converter_instance": mock_converter_instance,
        }


# =============================================================================
# Schema Factory Fixtures
# =============================================================================


@pytest.fixture
def slide_content_factory():
    """
    Factory for creating SlideContent objects.

    Usage:
        def test_something(slide_content_factory):
            slide = slide_content_factory(title="My Slide", bullet_points=["A", "B"])
    """
    from app.schemas import SlideContent

    def _create(
        layout_index: int = 1,
        title: str = "Test Slide",
        bullet_points: list[str] | None = None,
        image_url: str | None = None,
        speaker_notes: str | None = None,
    ):
        return SlideContent(
            layout_index=layout_index,
            title=title,
            bullet_points=bullet_points or ["Point 1", "Point 2"],
            image_url=image_url,
            speaker_notes=speaker_notes,
        )

    return _create
