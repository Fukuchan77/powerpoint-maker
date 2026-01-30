from functools import lru_cache

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.schemas import LayoutInfo, MasterInfo, PlaceholderInfo, TemplateAnalysisResult


@lru_cache(maxsize=32)
def _analyze_file_cached(file_path: str, template_id: str) -> TemplateAnalysisResult:
    """Analyze template file and return structure (Cached)"""
    # Note: lru_cache arguments must be hashable. Strings are hashable.
    prs = Presentation(file_path)
    masters = []

    for idx, master in enumerate(prs.slide_masters):
        layouts = []
        for layout_idx, layout in enumerate(master.slide_layouts):
            placeholders = []
            for shape in layout.placeholders:
                ph_format = shape.placeholder_format
                ph_type = ph_format.type

                # Determine acceptable content types
                accepts = []
                if shape.has_text_frame:
                    accepts.append("text")
                if ph_type == PP_PLACEHOLDER.PICTURE:
                    accepts.append("image")
                if ph_type == PP_PLACEHOLDER.TABLE:
                    accepts.append("table")
                if ph_type == PP_PLACEHOLDER.CHART:
                    accepts.append("chart")
                if ph_type == PP_PLACEHOLDER.OBJECT:
                    # objects can often hold text or images
                    accepts.append("text")
                    accepts.append("image")

                placeholders.append(
                    PlaceholderInfo(
                        idx=ph_format.idx,
                        name=shape.name,
                        type=str(ph_type),
                        width=shape.width,
                        height=shape.height,
                        left=shape.left,
                        top=shape.top,
                        accepts=accepts,
                    )
                )

            layouts.append(LayoutInfo(index=layout_idx, name=layout.name, placeholders=placeholders))

        masters.append(MasterInfo(index=idx, name=master.name or f"Master {idx}", layouts=layouts))

    return TemplateAnalysisResult(filename=file_path.split("/")[-1], template_id=template_id, masters=masters)


class TemplateAnalyzer:
    def analyze(self, file_path: str, template_id: str = "") -> TemplateAnalysisResult:
        return _analyze_file_cached(file_path, template_id)


class LayoutRegistry:
    _instance = None

    def __new__(cls):
        if cls._instance is None:
            cls._instance = super(LayoutRegistry, cls).__new__(cls)
            cls._instance._cache = {}
            cls._instance.analyzer = TemplateAnalyzer()
        return cls._instance

    def get_or_analyze(self, path: str, template_id: str) -> TemplateAnalysisResult:
        """Get analysis from cache or analyze file if not cached"""
        # Delegate to the cached function
        return self.analyzer.analyze(path, template_id)

    def clear(self):
        _analyze_file_cached.cache_clear()
