import os
from io import BytesIO
from typing import List, Optional

import requests
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

from app.schemas import ChartData, SlideContent


class SlidePopulator:
    def __init__(self, slide, strict=False):
        self.slide = slide
        self.strict = strict
        self.errors = []

    def safe_get_placeholder(self, idx: int, fallback_idx: int = None):
        """Get placeholder with fallback index support"""
        try:
            return self.slide.placeholders[idx]
        except KeyError:
            if fallback_idx is not None:
                try:
                    return self.slide.placeholders[fallback_idx]
                except KeyError:
                    pass

            # Try to find by index if idx lookup fails (pure fallback)
            for ph in self.slide.placeholders:
                if ph.placeholder_format.idx == idx:
                    return ph

            self.errors.append(f"No placeholder with idx={idx} (fallback={fallback_idx})")
            return None

    def validate_content_type(self, placeholder, content_type: str) -> bool:
        """Validate if placeholder supports content type"""
        if not placeholder:
            return False

        ph_type = placeholder.placeholder_format.type
        valid_mappings = {
            "text": [
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.OBJECT,
            ],
            "image": [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY],
            "chart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY],
            "table": [PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY],
        }
        return ph_type in valid_mappings.get(content_type, [])

    def _contains_japanese(self, text: str) -> bool:
        for char in text:
            if (
                "\u3000" <= char <= "\u303f"
                or "\u3040" <= char <= "\u309f"
                or "\u30a0" <= char <= "\u30ff"
                or "\u4e00" <= char <= "\u9faf"
                or "\uff00" <= char <= "\uffef"
            ):
                return True
        return False

    def set_japanese_font(self, run, font_name="Meiryo UI"):
        """Set East Asian font using XML manipulation"""
        try:
            rpr = run._r.get_or_add_rPr()
            # Check if a:ea already exists
            ea = rpr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(rpr, qn("a:ea"))
            ea.set("typeface", font_name)

            # Also set formatting for latin text just in case
            run.font.name = font_name
        except Exception as e:
            self.errors.append(f"Failed to set Japanese font: {e}")

    def replace_text_preserve_format(self, paragraph, new_text):
        """Replace text in first run only, preserving all formatting"""
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = new_text
            if self._contains_japanese(new_text):
                self.set_japanese_font(run)
            return

        # Keep only first run, preserve its formatting
        # Removing subsequent runs
        p = paragraph._p
        first_run_found = False
        for elm in list(p):
            if elm.tag.endswith("r") and not elm.tag.endswith("Pr"):  # is a run (and not pPr)
                if not first_run_found:
                    first_run_found = True
                    continue
                p.remove(elm)

        # Update first run
        if paragraph.runs:
            run = paragraph.runs[0]
            run.text = new_text
            if self._contains_japanese(new_text):
                self.set_japanese_font(run)

    def insert_picture_fit(self, placeholder, image_data: bytes):
        """Insert image fitted within placeholder, preserving aspect ratio"""
        try:
            im = Image.open(BytesIO(image_data))
            image_width, image_height = im.size

            picture = placeholder.insert_picture(BytesIO(image_data))

            # Reset cropping
            picture.crop_top = picture.crop_left = picture.crop_bottom = picture.crop_right = 0

            # Calculate aspect ratios
            placeholder_ratio = float(picture.width) / float(picture.height)
            image_ratio = float(image_width) / float(image_height)

            # Shrink to fit
            if placeholder_ratio > image_ratio:
                picture.width = int(image_ratio * picture.height)
            else:
                picture.height = int(picture.width / image_ratio)

            return picture
        except Exception as e:
            self.errors.append(f"Image processing failed: {e}")
            return None

    def set_theme_color(self, run, theme_color_name: str):
        """Set theme color on a run"""
        if not theme_color_name:
            return
        try:
            # Map string name to MSO_THEME_COLOR enum
            color_enum = getattr(MSO_THEME_COLOR, theme_color_name, MSO_THEME_COLOR.ACCENT_1)
            run.font.color.theme_color = color_enum
        except Exception as e:
            self.errors.append(f"Failed to set theme color {theme_color_name}: {e}")

    def populate_bullets(self, placeholder, bullets, theme_color=None):
        """Populate text frame with bullets (strings or BulletPoint objects)"""
        tf = placeholder.text_frame
        tf.clear()

        for item in bullets:
            p = tf.add_paragraph()

            if isinstance(item, str):
                text = item
                level = 0
            elif hasattr(item, "text"):  # BulletPoint object
                text = item.text
                level = item.level
            else:
                continue

            p.text = text
            p.level = level

            # Check Japanese
            is_japanese = self._contains_japanese(text)

            # Apply formatting to first run if exists
            if p.runs:
                run = p.runs[0]
                if is_japanese:
                    self.set_japanese_font(run)

                if theme_color:
                    self.set_theme_color(run, theme_color)

    def insert_chart(self, placeholder, chart_data: ChartData):
        """Insert a chart into the placeholder"""
        try:
            # Map chart type string to XL_CHART_TYPE enum
            chart_type_map = {
                "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
                "LINE": XL_CHART_TYPE.LINE,
                "PIE": XL_CHART_TYPE.PIE,
                "AREA": XL_CHART_TYPE.AREA,
                # Add more mappings as needed
            }

            # Default to COLUMN_CLUSTERED if type not found or invalid
            xl_chart_type = chart_type_map.get(chart_data.type.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED)

            # Create chart data object
            # Note: CategoryChartData works for Column, Bar, Line, Area.
            # Pie charts might need different handling if they don't support multiple series in same way,
            # but usually CategoryChartData is fine for simple Pie charts (one series).
            data = CategoryChartData()
            data.categories = chart_data.categories

            for series in chart_data.series:
                data.add_series(series.name, series.values)

            # Insert chart
            chart = placeholder.insert_chart(xl_chart_type, data).chart

            # Set title if supported
            try:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_data.title
            except Exception:
                pass  # Some chart types or layouts might not support title easily or raise error

        except Exception as e:
            self.errors.append(f"Failed to insert chart: {e}")


class PresentationGenerator:
    def _find_placeholder(self, slide, prefer_types: List[PP_PLACEHOLDER]) -> Optional[object]:
        """Find the best matching placeholder based on type priority"""
        for ph_type in prefer_types:
            for ph in slide.placeholders:
                if ph.placeholder_format.type == ph_type:
                    return ph
        return None

    def generate(self, template_path: str, slides: List[SlideContent], output_path: str) -> str:
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        prs = Presentation(template_path)

        # Clear existing slides
        xml_slides = prs.slides._sldIdLst
        slides_len = len(xml_slides)
        for i in range(slides_len - 1, -1, -1):
            del xml_slides[i]

        master = prs.slide_masters[0]  # Default to first master

        for slide_content in slides:
            # Layout selection
            if slide_content.layout_index >= len(master.slide_layouts):
                print(f"Warning: Layout index {slide_content.layout_index} out of range. Skipping slide.")
                continue

            layout = master.slide_layouts[slide_content.layout_index]
            slide = prs.slides.add_slide(layout)
            populator = SlidePopulator(slide)

            # 1. Handle Title
            if slide.shapes.title:
                populator.replace_text_preserve_format(slide.shapes.title.text_frame.paragraphs[0], slide_content.title)

            # 2. Handle Bullet Points (Body Text)
            # Priority: BODY > OBJECT
            bullets_to_use = slide_content.bullets or slide_content.bullet_points
            if bullets_to_use:
                body_placeholder = self._find_placeholder(slide, [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT])
                if body_placeholder and body_placeholder.has_text_frame:
                    populator.populate_bullets(body_placeholder, bullets_to_use, slide_content.theme_color)
                else:
                    print(f"Warning: No suitable placeholder found for text on slide '{slide_content.title}'")

            # 3. Handle Image
            if slide_content.image_url:
                # Priority: PICTURE > OBJECT > BODY
                pic_placeholder = self._find_placeholder(
                    slide, [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY]
                )

                if pic_placeholder:
                    try:
                        # Validate URL simple check
                        if not slide_content.image_url.startswith(("http://", "https://")):
                            print(f"Invalid image URL: {slide_content.image_url}")
                            continue

                        resp = requests.get(slide_content.image_url, timeout=10)  # 10s timeout
                        if resp.status_code == 200:
                            # If fallback to BODY, insert_picture might simply work if it's a placeholder?
                            # python-pptx placeholders usually have insert_picture method.
                            populator.insert_picture_fit(pic_placeholder, resp.content)
                        else:
                            print(f"Failed to fetch image: status {resp.status_code}")
                    except Exception as e:
                        print(f"Failed to fetch/insert image: {e}")
                else:
                    print(f"Warning: No suitable placeholder found for image on slide '{slide_content.title}'")

            # 4. Handle Chart
            if slide_content.chart:
                # Priority: CHART > OBJECT > BODY
                chart_placeholder = self._find_placeholder(
                    slide, [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY]
                )

                if chart_placeholder:
                    populator.insert_chart(chart_placeholder, slide_content.chart)
                else:
                    print(f"Warning: No suitable placeholder found for chart on slide '{slide_content.title}'")

        prs.save(output_path)
        return output_path
