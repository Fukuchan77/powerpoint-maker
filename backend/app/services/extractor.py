"""PPTX content extraction service.

Implements REQ-1.1.1~REQ-1.1.5, REQ-1.2.1, REQ-5.1
"""

import json
import uuid
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Optional

import structlog
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import EXTRACTED_IMAGE_EXPIRY_HOURS, EXTRACTED_IMAGES_DIR
from app.schemas import (
    AnalysisMode,
    BulletPoint,
    ChartSeries,
    ContentExtractionResult,
    ExtractedChart,
    ExtractedImage,
    ExtractedSlideContent,
)

logger = structlog.get_logger(__name__)


class ContentExtractor:
    """Extract content from PPTX files.

    Supports two modes:
    - CONTENT: Extract text, images, and charts
    - TEMPLATE: Extract layout information only (delegates to existing template analyzer)
    """

    def __init__(self, base_url: str):
        """Initialize the extractor.

        Args:
            base_url: Base URL for generating image access URLs
        """
        self.base_url = base_url.rstrip("/")
        self.logger = structlog.get_logger(__name__)

    async def extract(self, file_path: Path, mode: AnalysisMode) -> ContentExtractionResult:
        """Extract content from PPTX file [REQ-1.1.1~REQ-1.1.5]

        Args:
            file_path: Path to the PPTX file
            mode: Analysis mode (CONTENT or TEMPLATE)

        Returns:
            ContentExtractionResult with extracted slides, images, and warnings
        """
        import asyncio

        extraction_id = str(uuid.uuid4())
        extraction_dir = EXTRACTED_IMAGES_DIR / extraction_id
        extraction_dir.mkdir(parents=True, exist_ok=True)
        (extraction_dir / "images").mkdir(exist_ok=True)

        prs = Presentation(str(file_path))
        slides: list[ExtractedSlideContent] = []
        images: list[ExtractedImage] = []
        warnings: list[str] = []

        # Run slide extraction concurrently in thread pool to avoid blocking event loop
        # python-pptx is synchronous and CPU/IO bound
        tasks = []
        for slide_idx, slide in enumerate(prs.slides):
            task = asyncio.to_thread(
                self._extract_slide,
                slide,
                slide_idx,
                extraction_id,
                extraction_dir,
                mode,
            )
            tasks.append(task)

        if tasks:
            results = await asyncio.gather(*tasks)

            # Aggregate results properly sorted by slide index (gather maintains order)
            for slide_content, slide_images, slide_warnings in results:
                slides.append(slide_content)
                images.extend(slide_images)
                warnings.extend(slide_warnings)

        expires_at = datetime.now(UTC).replace(tzinfo=None) + timedelta(hours=EXTRACTED_IMAGE_EXPIRY_HOURS)

        # Save metadata for cleanup service
        metadata = {
            "expires_at": expires_at.isoformat(),
            "image_count": len(images),
            "filename": file_path.name,
        }
        (extraction_dir / "metadata.json").write_text(json.dumps(metadata))

        self.logger.info(
            "content_extracted",
            extraction_id=extraction_id,
            slide_count=len(slides),
            image_count=len(images),
            warning_count=len(warnings),
        )

        return ContentExtractionResult(
            extraction_id=extraction_id,
            filename=file_path.name,
            expires_at=expires_at.isoformat(),
            slides=slides,
            images=images,
            warnings=warnings,
        )

    def _extract_slide(
        self,
        slide,
        slide_idx: int,
        extraction_id: str,
        extraction_dir: Path,
        mode: AnalysisMode,
    ) -> tuple[ExtractedSlideContent, list[ExtractedImage], list[str]]:
        """Extract content from a single slide.

        Returns:
            Tuple of (slide_content, list_of_images, list_of_warnings)
        """
        # Local containers for this slide's artifacts
        images: list[ExtractedImage] = []
        warnings: list[str] = []

        # Get layout index
        layout_idx = 0
        try:
            layout_idx = slide.slide_layout.slide_master.slide_layouts.index(slide.slide_layout)
        except (ValueError, AttributeError):
            pass

        title: Optional[str] = None
        body_text: list[str] = []
        bullet_points: list[BulletPoint] = []
        image_refs: list[str] = []
        chart: Optional[ExtractedChart] = None

        for shape in slide.shapes:
            try:
                # Extract title
                if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type == 1:  # TITLE
                        if shape.has_text_frame:
                            title = shape.text_frame.text.strip()
                        continue

                # Skip content extraction in template mode [REQ-1.2.1]
                if mode == AnalysisMode.TEMPLATE:
                    continue

                # Extract text
                if shape.has_text_frame:
                    texts = self._extract_text_from_shape(shape)
                    body_text.extend(texts)
                    for text in texts:
                        bullet_points.append(BulletPoint(text=text, level=0))

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    extracted = self._extract_image(shape, extraction_id, slide_idx, extraction_dir)
                    if extracted:
                        images.append(extracted)
                        image_refs.append(extracted.id)

                # Extract charts
                if shape.has_chart:
                    extracted_chart = self._extract_chart(shape, slide_idx)
                    if extracted_chart:
                        chart = extracted_chart

            except Exception as e:
                warning = self._handle_unsupported_shape(shape, slide_idx, str(e))
                warnings.append(warning)
                self.logger.warning("shape_extraction_failed", slide=slide_idx, error=str(e))

        slide_content = ExtractedSlideContent(
            slide_index=slide_idx,
            layout_index=layout_idx,
            title=title,
            body_text=body_text,
            bullet_points=bullet_points,
            image_refs=image_refs,
            chart=chart,
        )

        return slide_content, images, warnings

    def _extract_text_from_shape(self, shape) -> list[str]:
        """Extract text from shape [REQ-1.1.1]

        Args:
            shape: pptx shape object

        Returns:
            List of text strings from paragraphs
        """
        texts = []
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
        return texts

    def _extract_image(
        self,
        shape,
        extraction_id: str,
        slide_idx: int,
        extraction_dir: Path,
    ) -> Optional[ExtractedImage]:
        """Extract and save image [REQ-1.1.2]

        Args:
            shape: Picture shape object
            extraction_id: Unique extraction session ID
            slide_idx: Slide index
            extraction_dir: Directory to save images

        Returns:
            ExtractedImage if successful, None otherwise
        """
        try:
            image = shape.image
            image_id = str(uuid.uuid4())
            ext = image.ext or "png"
            filename = f"{image_id}.{ext}"
            image_path = extraction_dir / "images" / filename
            image_path.write_bytes(image.blob)

            return ExtractedImage(
                id=image_id,
                filename=shape.name or filename,
                url=f"{self.base_url}/api/extracted-images/{extraction_id}/{image_id}",
                slide_index=slide_idx,
                content_type=image.content_type or f"image/{ext}",
            )
        except AttributeError as e:
            self.logger.warning("image_extraction_failed", slide=slide_idx, error=str(e))
            return None

    def _extract_chart(self, shape, slide_idx: int) -> Optional[ExtractedChart]:
        """Extract chart data [REQ-1.1.4]

        Args:
            shape: Chart shape object
            slide_idx: Slide index

        Returns:
            ExtractedChart if successful, None otherwise
        """
        try:
            chart = shape.chart
            plot = chart.plots[0]

            categories = []
            if plot.categories:
                categories = [str(c) for c in plot.categories]

            series = []
            for s in plot.series:
                series.append(ChartSeries(name=s.name or "", values=list(s.values)))

            return ExtractedChart(
                slide_index=slide_idx,
                chart_type=str(chart.chart_type),
                categories=categories,
                series=series,
            )
        except Exception as e:
            self.logger.warning("chart_extraction_failed", slide=slide_idx, error=str(e))
            return None

    def _handle_unsupported_shape(self, shape, slide_idx: int, error: str) -> str:
        """Generate warning for unsupported elements [REQ-5.1]

        Args:
            shape: Shape that couldn't be processed
            slide_idx: Slide index
            error: Error message

        Returns:
            Warning message string
        """
        shape_type = getattr(shape, "shape_type", "unknown")
        return f"Slide {slide_idx + 1}: Unsupported or failed element '{shape_type}' - {error}"
